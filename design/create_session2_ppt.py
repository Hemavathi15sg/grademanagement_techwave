from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
TITLE_COLOR = RGBColor(31, 78, 121)  # Dark Blue
ACCENT_COLOR = RGBColor(255, 87, 34)  # Orange
TEXT_COLOR = RGBColor(51, 51, 51)  # Dark Gray
LIGHT_BG = RGBColor(245, 245, 245)  # Light Gray

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = ACCENT_COLOR
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(12)
        p.level = 0
    
    return slide

def add_two_column_slide(prs, title, left_items, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(10)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(10)
    
    return slide

# Slide 1: Title
slide1 = add_title_slide(prs, 
    "SESSION 2: Refactoring & QA Expansion",
    "🔧 From Technical Debt to Enterprise Architecture")

slide1.notes_slide.notes_text_frame.text = """
ENERGY: Start with HIGH energy! Session 2 is about REALITY MODE.

OPENING (30 seconds):
"Welcome back! Session 1 was the dream - building NEW features from scratch. Clean slate, perfect architecture, greenfield development.

But here's reality: Most of your time ISN'T spent on greenfield. You inherit code. You maintain legacy systems. You fix someone else's mess.

Session 2 is REALITY MODE. Today, we're taking MESSY existing code - the kind you actually deal with every day - and transforming it into clean, enterprise-grade architecture using AI.

Technical debt doesn't have to compound. With GitHub Copilot, paying it down is FAST."
"""

# Slide 2: Session Overview
slide2 = add_content_slide(prs,
    "Session 2 Overview",
    [
        "⏱️ Duration: 25 minutes",
        "",
        "🎯 Objective:",
        "Transform existing messy code with technical debt into",
        "clean, testable, enterprise-ready architecture",
        "",
        "👥 Audience:",
        "Developers, QA Engineers, Tech Leads",
        "",
        "📍 Starting Point:",
        "Existing student/grade handlers with global state,",
        "mixed concerns, and tight coupling",
        "",
        "🏆 Ending Point:",
        "Repository pattern, comprehensive test infrastructure,",
        "87% test coverage"
    ])

slide2.notes_slide.notes_text_frame.text = """
SPEAKER NOTES:

Set the context clearly:
"We're NOT in the demo/ folder from Session 1. This is a DIFFERENT codebase - existing code with real technical debt. This is authentic brownfield refactoring.

The code we're starting with has THREE deadly sins:
1. Global state - shared maps, impossible to test in parallel
2. Mixed responsibilities - HTTP logic AND database logic in same function  
3. Tight coupling - can't swap storage without rewriting everything

This isn't a teaching example. This is REAL technical debt that real teams face every day."

Pause for effect after describing the technical debt. Let it sink in.
"""

# Slide 3: How GitHub Copilot Helps in Session 2
slide3 = add_two_column_slide(prs,
    "How GitHub Copilot Helps in Session 2",
    [
        "🤖 Three Agent Types:",
        "",
        "1️⃣ Copilot Coding Agent",
        "• Multi-file refactoring",
        "• Architectural transformation",
        "• Coordinated changes across 5+ files",
        "• Creates PRs automatically",
        "",
        "2️⃣ Local Agent",
        "• Interactive code generation",
        "• Test infrastructure creation",
        "• Iterative refinement",
        "• Real-time feedback",
        "",
        "3️⃣ Background Agent",
        "• Parallel documentation",
        "• Pattern generation",
        "• Non-blocking tasks"
    ],
    [
        "💡 AI-Powered Capabilities:",
        "",
        "✅ Extract repository pattern from",
        "   messy global state automatically",
        "",
        "✅ Generate 500+ lines of mock code",
        "   with gomock in seconds",
        "",
        "✅ Create test data factories with",
        "   builder pattern instantly",
        "",
        "✅ Write comprehensive integration",
        "   tests covering all scenarios",
        "",
        "✅ Maintain 100% API compatibility",
        "   during refactoring",
        "",
        "⚡ Result:",
        "4.2x faster refactoring cycles",
        "87% test coverage achieved",
        "Zero breaking changes"
    ])

slide3.notes_slide.notes_text_frame.text = """
SPEAKER NOTES:

"The key difference in Session 2: We're using MULTIPLE agent types strategically.

COPILOT CODING AGENT: When refactoring spans 5+ files, we need coordination. The agent analyzes dependencies, plans changes across the entire codebase, and executes them in one coordinated transformation. Manual refactoring? One mistake breaks everything. Agent? It handles all the complexity.

LOCAL AGENT: For generating test infrastructure - mocks and factories - we use Local Agent. Quick generation, iterative refinement. We can see the code immediately and adjust.

The magic: GitHub Copilot doesn't just generate NEW code fast. It helps TRANSFORM EXISTING code fast. That's the breakthrough.

Traditional refactoring: Nerve-wracking, time-consuming, risky.
AI-assisted refactoring: Confident, fast, safe.

This session proves: Technical debt is no longer expensive to fix. With AI, continuous improvement is economically viable."

TIMING: Spend 1-2 minutes on this slide. It's the foundation for understanding the demos.
"""

# Slide 4: Session 2 Demo Topics
slide4 = add_content_slide(prs,
    "Session 2 Demo Topics",
    [
        "🏗️ Act 1: Repository Refactoring Agent (8 minutes)",
        "   • Extract repository pattern from messy handlers",
        "   • Remove global state and inject dependencies",
        "   • Surgical multi-file transformation",
        "   • Time saved: 15 min → 8 min (1.9x faster)",
        "",
        "🧪 Act 2: Mock + Factory Generation Agent (10 minutes)",
        "   • Auto-generate gomock mocks for repositories",
        "   • Create test data factories with builder pattern",
        "   • 500+ lines of test infrastructure in minutes",
        "   • Time saved: 45 min → 10 min (4.5x faster)",
        "",
        "🎯 Act 3: Integration Test Expansion Agent (7 minutes)",
        "   • Comprehensive CRUD workflow tests",
        "   • Concurrent testing for thread safety",
        "   • Coverage: 23% → 87% in 7 minutes",
        "   • Time saved: 45 min → 7 min (6.4x faster)",
        "",
        "💰 Total Time Saved: 105 minutes → 25 minutes (4.2x speedup)"
    ])

slide4.notes_slide.notes_text_frame.text = """
SPEAKER NOTES:

"Three acts, three transformations, one mission: Turn technical debt into enterprise architecture.

ACT 1 - REPOSITORY REFACTORING: This is the foundation. We're taking handlers that directly manipulate global variables and extracting a clean repository pattern. The agent coordinates changes across 5 files simultaneously. Traditional approach? 15 minutes of nerve-wracking surgery. AI approach? 8 minutes of confident delegation.

ACT 2 - TEST INFRASTRUCTURE: Clean architecture means nothing without tests. We auto-generate mocks using gomock and create beautiful test data factories with builder pattern. Writing these manually? 45 minutes of boilerplate. With AI? 10 minutes of generation.

ACT 3 - TEST EXPANSION: Now we USE our infrastructure. We write comprehensive integration tests covering happy paths, error scenarios, concurrent requests, edge cases. The result? Coverage jumps from 23% to 87% in 7 minutes.

The numbers tell the story: 105 minutes of traditional work becomes 25 minutes with AI. But it's not just speed - it's CONFIDENCE. Every step is tested. Every change is verified. This is sustainable velocity."

TRANSITION: "Let's dive into Act 1 and watch AI perform architectural surgery on production code..."
"""

# Slide 5: Act 1 - Repository Refactoring
slide5 = add_two_column_slide(prs,
    "Act 1: Repository Refactoring Agent",
    [
        "❌ BEFORE: Technical Debt",
        "",
        "• Global variables polluting scope:",
        "  var students = make(map[int]...)",
        "  var studentsMu sync.RWMutex",
        "",
        "• Mixed responsibilities:",
        "  HTTP + Database logic together",
        "",
        "• Tight coupling:",
        "  Can't swap storage",
        "  Can't test in isolation",
        "",
        "• Impossible to mock",
        "",
        "• Race conditions possible",
        "",
        "🚨 Result:",
        "Unmaintainable, untestable,",
        "inflexible architecture"
    ],
    [
        "✅ AFTER: Clean Architecture",
        "",
        "• Repository interfaces:",
        "  StudentRepository interface",
        "  with Create, Get, Update, Delete",
        "",
        "• Dependency injection:",
        "  type StudentHandler struct {",
        "    Repo repository.StudentRepository",
        "  }",
        "",
        "• Encapsulated state:",
        "  All storage inside repository",
        "",
        "• Thread-safe implementations",
        "",
        "🏆 Result:",
        "Testable, maintainable,",
        "flexible, production-ready",
        "",
        "⏱️ Time: 8 minutes via Copilot Agent"
    ])

slide5.notes_slide.notes_text_frame.text = """
SPEAKER NOTES:

"Let me show you the transformation visually.

BEFORE: This is what we inherited. Global variables everywhere. Look at line 14-17 in the handler: var students, var studentsMu, var nextStudentID. All polluting package scope. 

In CreateStudent function: Seven responsibilities in one place - HTTP decoding, validation, mutex locking, ID generation, storage mutation, serialization. This is a MAINTENANCE NIGHTMARE.

AFTER: Look at this beauty. Clean repository interface defining the contract. Handler depends on the INTERFACE, not concrete implementation. All state encapsulated inside repository structs. Constructor injection for testability.

The handler is now PURE - it only handles HTTP and delegates storage. Single responsibility principle achieved.

THE DEMO: We'll create Jira story TEC-11, delegate to Copilot Coding Agent with detailed architectural requirements, watch it coordinate changes across 5 files, review the generated PR, and then VERIFY with live API testing.

The magic moment: When we run curl commands and the refactored API works IDENTICALLY to before. Zero breaking changes. Clean architecture achieved. That's professional refactoring."

DEMO CHECKPOINT: Show actual code comparison on screen if possible. Visual impact is powerful.
"""

# Slide 6: Act 2 - Mock & Factory Generation
slide6 = add_two_column_slide(prs,
    "Act 2: Mock + Factory Generation Agent",
    [
        "🎭 GOMOCK GENERATION:",
        "",
        "• Auto-generate from interfaces:",
        "  //go:generate mockgen...",
        "",
        "• MockStudentRepository",
        "  with expectation setting",
        "",
        "• MockGradeRepository",
        "  with call verification",
        "",
        "• 500+ lines generated",
        "  in seconds",
        "",
        "• Full EXPECT() syntax support",
        "",
        "⏱️ Manual work: 30 minutes",
        "⚡ AI generation: 3 minutes"
    ],
    [
        "🏭 FACTORY PATTERN:",
        "",
        "• Builder pattern with fluent API:",
        "  NewStudentBuilder()",
        "  .WithName('Alice')",
        "  .WithGrade('A')",
        "  .Build()",
        "",
        "• Sensible defaults included",
        "",
        "• Preset configurations:",
        "  ValidStudent()",
        "  HighAchiever()",
        "  FailingStudent()",
        "",
        "• BuildN(n) for bulk generation",
        "",
        "⏱️ Manual work: 15 minutes",
        "⚡ AI generation: 2 minutes"
    ])

slide6.notes_slide.notes_text_frame.text = """
SPEAKER NOTES:

"Act 2 is about TEST INFRASTRUCTURE. We have clean architecture, now we need tools to test it.

GOMOCK GENERATION: The agent adds go:generate directives to our repository interfaces and runs mockgen. Result? MockStudentRepository with full expectation syntax. We can now write unit tests like:

  mock.EXPECT().Create(gomock.Any()).Return(expectedStudent)
  
This is ISOLATION TESTING. We test the handler without touching any real storage. Fast tests. Reliable tests.

FACTORY PATTERN: Creating test data manually is PAINFUL. Every test needs:
  student := models.Student{
    ID: 1,
    Name: 'Alice Johnson',
    Email: 'alice@university.edu',
    Grade: 'A',
  }

That's 7 lines per test. With factory:
  student := ValidStudent().Build()

That's 1 line. 80% less test setup code.

THE DEMO: We'll create Jira TEC-7, use Local Agent to generate both mocks and factories, show example tests using them, and prove how much cleaner test code becomes.

The impact: Test setup goes from 'painful boilerplate' to 'one line of intention'. Tests become DOCUMENTATION."

ENERGY: Show excitement about how much easier testing becomes. This is about developer happiness.
"""

# Slide 7: Act 3 - Integration Test Expansion
slide7 = add_content_slide(prs,
    "Act 3: Integration Test Expansion Agent",
    [
        "📊 COVERAGE EXPLOSION:",
        "   Before: 23% coverage    →    After: 87% coverage",
        "",
        "🧪 COMPREHENSIVE TEST SUITE:",
        "",
        "• TestStudentCRUDWorkflow - Full lifecycle testing",
        "  Create → Read → Update → Delete → Verify 404",
        "",
        "• TestStudentConcurrentCreation - Thread safety validation",
        "  100 parallel requests with unique ID verification",
        "",
        "• TestStudentErrorScenarios - Sad path coverage",
        "  404 for missing resources, 400 for invalid input",
        "",
        "• TestStudentListPagination - Edge cases",
        "  Empty lists, bulk operations, boundary conditions",
        "",
        "⚡ Execution Speed: 134ms for entire suite",
        "🎯 Result: Tests serve as living documentation",
        "⏱️ Time Saved: 45 minutes → 7 minutes (6.4x faster)"
    ])

slide7.notes_slide.notes_text_frame.text = """
SPEAKER NOTES:

"Act 3 is where it all comes together. We have clean architecture. We have test infrastructure. Now we USE it to achieve enterprise-grade coverage.

THE COVERAGE JUMP: From 23% to 87% in 7 minutes. That's not just a number - that's CONFIDENCE.

COMPREHENSIVE TESTS:
- CRUD Workflow tests the happy path end-to-end
- Concurrent Creation proves our thread safety actually works under load
- Error Scenarios test that failures fail gracefully
- Edge Cases catch the bugs that slip through

THE MAGIC: These tests read like SPECIFICATIONS. Look at TestStudentCRUDWorkflow - you can read the test and understand exactly how the API should behave. That's documentation that never goes stale because it's EXECUTED.

THE DEMO: We'll create Jira TEC-12, use Local Agent to generate the test suite, run the tests live showing them all passing, and watch coverage jump in real-time.

CELEBRATION MOMENT: When you run 'go test -cover' and see '87.3% coverage' - that's your victory moment. Pause and let the audience appreciate that number.

The real win: These tests execute in 134ms. Fast feedback loops. You can run them every single commit."

PAUSE: Before moving to results, check audience engagement. This is a good moment for a quick question.
"""

# Slide 8: Session 2 Results & Impact
slide8 = add_two_column_slide(prs,
    "Session 2 Results & Impact",
    [
        "📈 QUANTITATIVE RESULTS:",
        "",
        "⏱️ Time Investment:",
        "   25 minutes guided refactoring",
        "",
        "💰 Time Saved:",
        "   Act 1: 15 min → 8 min",
        "   Act 2: 45 min → 10 min",
        "   Act 3: 45 min → 7 min",
        "   Total: 105 min → 25 min",
        "   Speedup: 4.2x",
        "",
        "📊 Quality Metrics:",
        "   Coverage: 23% → 87%",
        "   Test execution: 134ms",
        "   Files refactored: 5",
        "   Lines generated: 500+",
        "   Breaking changes: 0"
    ],
    [
        "🎯 QUALITATIVE IMPACT:",
        "",
        "✅ Architecture:",
        "   Clean repository pattern",
        "   SOLID principles followed",
        "   Dependency injection",
        "",
        "✅ Testability:",
        "   Can mock any component",
        "   Fast isolated unit tests",
        "   Comprehensive integration tests",
        "",
        "✅ Maintainability:",
        "   Single responsibility",
        "   Clear separation of concerns",
        "   Tests as documentation",
        "",
        "✅ Flexibility:",
        "   Can swap storage (PostgreSQL,",
        "   MongoDB, Redis) without",
        "   touching handlers"
    ])

slide8.notes_slide.notes_text_frame.text = """
SPEAKER NOTES:

"Let's talk about the IMPACT. Not just speed - TRANSFORMATION.

QUANTITATIVE RESULTS: 4.2x speedup is impressive. But look deeper:
- Zero breaking changes means we can deploy immediately
- 87% coverage means we catch bugs before production
- 134ms test execution means fast feedback loops
- 500+ lines generated means less manual tedious work

QUALITATIVE IMPACT: This is where it gets real:

FOR DEVELOPERS: You're not afraid to change code anymore. Tests protect you. Architecture is clean. New features are easy to add.

FOR QA: Your tests are comprehensive and readable. They document expected behavior. Coverage gives you confidence.

FOR TECH LEADS: Technical debt is eliminated without slowing velocity. You can onboard new developers 3x faster because code is clean and documented.

FOR ENGINEERING LEADERS: This proves ROI on AI tooling. 4.2x faster refactoring means you can afford continuous improvement. Quality and velocity are NOT opposites.

THE BIGGER PICTURE: Session 1 showed speed. Session 2 shows SUSTAINABILITY. With AI, you don't have to choose between shipping fast and maintaining quality. You can do BOTH.

This is the evolution of software development: Continuous improvement is no longer expensive."

TRANSITION: "So what are the key lessons we take away from this?"
"""

# Slide 9: Key Takeaways
slide9 = add_content_slide(prs,
    "Key Takeaways from Session 2",
    [
        "1️⃣ REFACTOR EARLY WITH AI ASSISTANCE",
        "   Don't wait for debt to compound | AI makes refactoring economically viable",
        "",
        "2️⃣ INTERFACES UNLOCK TESTABILITY",
        "   Repository pattern + dependency injection = mockable code",
        "   Isolation testing is fast testing | Testable architecture enables rapid development",
        "",
        "3️⃣ FACTORIES MAKE TESTS MAINTAINABLE",
        "   Test code IS documentation | Builder pattern eliminates duplication",
        "   80% less test setup = more time testing",
        "",
        "4️⃣ EVOLUTION > REVOLUTION",
        "   Continuous improvement with AI | Maintain velocity while fixing debt",
        "   Quality and speed are NOT opposites",
        "",
        "5️⃣ AI TRANSFORMS EXISTING CODE, NOT JUST NEW CODE",
        "   The breakthrough: AI helps you MAINTAIN code faster",
        "   Architecture evolves, never rots | Technical debt becomes manageable"
    ])

slide9.notes_slide.notes_text_frame.text = """
SPEAKER NOTES:

"Five core lessons from Session 2. These are principles that change how you think about maintenance:

1. REFACTOR EARLY: The old way was 'ship fast, refactor later' - but later never came. With AI, refactoring is FAST. You can afford to do it continuously. Don't wait for debt to compound. Fix it early while it's small.

2. INTERFACES UNLOCK TESTABILITY: The repository pattern isn't just about clean code - it's about SPEED. When you can test in isolation with mocks, tests run in milliseconds instead of seconds. Fast tests mean fast feedback. Fast feedback means high velocity.

3. FACTORIES MAKE TESTS MAINTAINABLE: Here's the dirty secret: Most developers hate writing tests because test setup is PAINFUL. Factories fix that. When creating test data is easy, testing becomes pleasant. When testing is pleasant, coverage goes up.

4. EVOLUTION > REVOLUTION: Stop thinking in terms of 'big rewrites'. Think continuous evolution. With AI, small improvements are economically viable. Ship features fast. Refactor fast. Test fast. Evolve continuously.

5. AI TRANSFORMS EXISTING CODE: This is the breakthrough most people miss. AI isn't just for greenfield. It's for BROWNFIELD. Your legacy code can become clean code. Your technical debt can be paid down. That's the real revolution.

CALL TO ACTION: 'This week, pick ONE messy handler. Spend 30 minutes refactoring it with Copilot Agent. I guarantee you'll see the power of AI-assisted evolution.'"
"""

# Slide 10: What's Next
slide10 = add_content_slide(prs,
    "What's Next: Sessions 3 & 4",
    [
        "📅 SESSION 3: Security & DevOps Automation",
        "   • Security scanning with AI-powered remediation",
        "   • Automated CI/CD pipeline generation",
        "   • Infrastructure as Code with AI assistance",
        "   • Production deployment strategies",
        "",
        "📅 SESSION 4: Advanced Workflows & Multi-Agent Orchestration",
        "   • Complex feature development with multiple agents",
        "   • Code review automation and quality gates",
        "   • Production monitoring and incident response",
        "   • Enterprise scaling patterns",
        "",
        "🎯 YOUR HOMEWORK THIS WEEK:",
        "   Day 1: Pick one messy handler with global state",
        "   Day 2: Delegate repository extraction to Copilot Agent",
        "   Day 3: Generate mocks with gomock",
        "   Day 4: Create test factories for your domain",
        "   Day 5: Write 10 integration tests using new infrastructure",
        "",
        "By Friday: Clean architecture + test coverage you can trust!"
    ])

slide10.notes_slide.notes_text_frame.text = """
SPEAKER NOTES:

"We've covered building fast (Session 1) and maintaining fast (Session 2). But what about DEPLOYING fast and SCALING fast? That's Sessions 3 and 4.

SESSION 3 - SECURITY & DEVOPS: AI doesn't just write code. It can scan for vulnerabilities, generate CI/CD pipelines, write infrastructure as code, and automate deployment. We'll show how to go from code to production with AI assistance at every step.

SESSION 4 - ADVANCED WORKFLOWS: The final session brings it all together. Multi-agent orchestration for complex features. Automated code review. Production monitoring. This is the full end-to-end enterprise workflow.

YOUR HOMEWORK: Don't just watch - DO. Here's your 5-day challenge:
- Day 1: Identify the problem (messy handler)
- Day 2: Fix the architecture (repository pattern)
- Day 3: Add mocking capability (gomock)
- Day 4: Make testing easy (factories)
- Day 5: Prove it works (integration tests)

By Friday, you'll have transformed ONE piece of technical debt into clean, tested, production-ready code. That's your proof of concept.

NEXT WEEK: Do it again with another handler. Before you know it, your entire codebase is clean.

This is continuous evolution. This is sustainable velocity. This is the future of software development."

END ON HIGH NOTE: "Questions? Who's ready to refactor their first handler this week?"
"""

# Slide 11: Questions & Discussion
slide11 = add_title_slide(prs,
    "Questions & Discussion",
    "Share your technical debt challenges!")

slide11.notes_slide.notes_text_frame.text = """
SPEAKER NOTES:

OPEN THE FLOOR:

Good prompting questions:
1. "Who has handlers with global state in their codebase?" (Show of hands)
2. "What's the BIGGEST piece of technical debt you're carrying?" (Open discussion)
3. "What stops you from refactoring today?" (Identify barriers)
4. "How many of you have test coverage below 50%?" (Pain point validation)

HANDLE COMMON QUESTIONS:

Q: "What if our codebase is way messier than your example?"
A: "Start smaller. One function at a time. AI helps with any size refactoring. The messier the code, the MORE value you get from AI assistance."

Q: "How do we convince management to allocate time for refactoring?"
A: "Show them the numbers. 4.2x faster means you can refactor WITHOUT slowing feature delivery. It's not either/or anymore - it's both."

Q: "What if tests fail after refactoring?"
A: "That's the beauty of comprehensive tests - they tell you immediately. And because you can revert via git, there's no risk. Safety net exists."

Q: "Do we need to learn gomock and all these patterns first?"
A: "No! That's the point. AI generates the patterns. You learn by seeing the generated code. Reverse learning - output before theory."

CLOSING ENERGY:
"Thank you! Remember: Technical debt is no longer your enemy. With AI, paying it down is FAST. See you in Session 3 for Security & DevOps! 🚀"
"""

# Save presentation
prs.save('SESSION2_Refactoring_QA_Expansion.pptx')
print("✅ PowerPoint created successfully: SESSION2_Refactoring_QA_Expansion.pptx")
print("📊 11 slides generated with comprehensive speaker notes")
print("🎨 Professional design with custom color scheme")
print("🎤 Detailed speaker notes for confident delivery")
