"""
Generate Session 1 PowerPoint Presentation
GitHub Copilot Agent Mode, Boilerplate & Documentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
COLOR_PRIMARY = RGBColor(26, 35, 126)  # Deep Blue
COLOR_SECONDARY = RGBColor(124, 77, 255)  # Electric Purple
COLOR_ACCENT = RGBColor(0, 229, 255)  # Bright Cyan
COLOR_SUCCESS = RGBColor(0, 230, 118)  # Vibrant Green
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_DARK = RGBColor(30, 30, 30)

def add_title_slide(prs, title, subtitle):
    """Add title slide with gradient background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background effect (simulated with shapes)
    background = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_PRIMARY
    background.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2), Inches(8), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = COLOR_ACCENT
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_lines, notes=""):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(245, 245, 250)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i > 0:
            content_frame.add_paragraph()
        p = content_frame.paragraphs[i]
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_DARK
        p.space_before = Pt(12)
        p.level = 0
        if line.startswith('  '):
            p.level = 1
            p.text = line.strip()
    
    # Add speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide

def add_stats_slide(prs, title, stats_data, notes=""):
    """Add slide with statistics and numbers"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(245, 245, 250)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    
    # Stats boxes
    y_position = 1.5
    for stat in stats_data:
        stat_box = slide.shapes.add_shape(
            1, Inches(1), Inches(y_position), Inches(8), Inches(1.2)
        )
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = COLOR_WHITE
        stat_box.line.color.rgb = COLOR_SECONDARY
        stat_box.line.width = Pt(2)
        
        stat_frame = stat_box.text_frame
        stat_frame.text = stat
        stat_para = stat_frame.paragraphs[0]
        stat_para.font.size = Pt(22)
        stat_para.font.bold = True
        stat_para.font.color.rgb = COLOR_PRIMARY
        stat_para.alignment = PP_ALIGN.CENTER
        stat_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        y_position += 1.4
    
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide

# SLIDE 1: Title Slide
add_title_slide(
    prs,
    "🚀 FROM ZERO TO PRODUCTION\nIN 30 MINUTES",
    "GitHub Copilot Agent Mode • Boilerplate & Documentation"
)

# SLIDE 2: The Developer's Struggle
add_content_slide(
    prs,
    "⏰ THE DEVELOPER'S DAILY STRUGGLE",
    [
        "YOUR TYPICAL SPRINT WEEK:",
        "",
        "🏗️  CRUD Boilerplate → 16 hours",
        "⚡  Performance Tuning → 8 hours",
        "🧪  Writing Tests → 12 hours",
        "📝  Documentation → 4 hours",
        "❓  Meetings & Code Review → 10 hours",
        "",
        "💡 ACTUAL INNOVATION → 2 hours",
        "",
        "52 hours of work, only 2 hours of creativity"
    ],
    "Last week, I tracked my team's hours. 16 hours on CRUD boilerplate. 8 hours tweaking performance. 12 hours writing tests. 4 hours on documentation. Just 2 hours on actual innovation. This is the productivity crisis nobody talks about. What if AI could handle the patterns while we focus on creativity?"
)

# SLIDE 3: Session Agenda
add_content_slide(
    prs,
    "📋 SESSION 1 AGENDA",
    [
        "🎯 OBJECTIVE: Build Production-Ready API in 30 Minutes",
        "",
        "ACT 1: CRUD Boilerplate Agent (8 min)",
        "  → Complete enrollment API generation",
        "",
        "ACT 2: Performance Agent (6 min)",
        "  → Redis caching & optimization",
        "",
        "ACT 3: QA Contract Agent (8 min)",
        "  → OpenAPI specs & testing",
        "",
        "ACT 4: Documentation Agent (5 min)",
        "  → Comprehensive docs generation",
        "",
        "💰 EXPECTED ROI: 10 hours saved, $750+ value"
    ],
    "We're building a complete grade management system through four acts. Act 1: CRUD API in 8 minutes (normally 4 hours). Act 2: Redis caching in 6 minutes (normally 2 hours). Act 3: OpenAPI specs and tests in 8 minutes (normally 3 hours). Act 4: Documentation in 5 minutes (normally 1 hour). Total: 10 hours saved, over $750 in value."
)

# SLIDE 4: Agent Types
add_content_slide(
    prs,
    "🤖 MEET YOUR AI WORKFORCE",
    [
        "🌥️  CLOUD CODING AGENT",
        "  ✓ Complex feature development",
        "  ✓ Multi-file code generation",
        "  ✓ Automatic PR creation",
        "  Best for: Complete features, new architecture",
        "",
        "💻 LOCAL AGENT",
        "  ✓ Interactive code suggestions",
        "  ✓ Rapid iteration & debugging",
        "  ✓ Real-time feedback loops",
        "  Best for: Optimization, testing, quick fixes",
        "",
        "🔄 BACKGROUND AGENT",
        "  ✓ Parallel task execution",
        "  ✓ Non-blocking workflows",
        "  ✓ Independent operations",
        "  Best for: Documentation, formatting, cleanup"
    ],
    "GitHub Copilot has three specialized agents. Cloud Coding Agent: your senior developer for complex features, creates PRs automatically. Local Agent: your pair programmer for real-time suggestions and debugging. Background Agent: your parallel worker for independent tasks like docs. The magic happens when you use all three together."
)

# SLIDE 5: Delegation Strategy
add_content_slide(
    prs,
    "🎯 THE AI DELEGATION STRATEGY",
    [
        "STRATEGIC AGENT ASSIGNMENT:",
        "",
        "New Features → 🌥️ Cloud Coding Agent",
        "  Complex, multi-file, full context needed",
        "",
        "Performance Tuning → 💻 Local Agent",
        "  Iterative testing, real-time feedback",
        "",
        "Quality Assurance → 💻 Local Agent",
        "  Validation needed, immediate results",
        "",
        "Documentation → 🔄 Background Agent",
        "  Independent work, parallel execution",
        "",
        "🎪 THE PRINCIPLE:",
        "Humans define WHAT → AI determines HOW"
    ],
    "Strategic delegation: Cloud Agent for complex features needing full context. Local Agent for iterative work like performance tuning and testing. Background Agent for independent tasks like documentation. The principle: humans define WHAT to build, AI determines HOW to build it. This is maximum productivity."
)

# SLIDE 6: Act 1 - CRUD
add_stats_slide(
    prs,
    "🏗️ ACT 1: CRUD BOILERPLATE AGENT",
    [
        "Traditional: 4 hours of manual coding",
        "AI Delegation: 8 minutes with Cloud Agent",
        "⚡ 30x FASTER",
        "",
        "DELIVERED: Complete enrollment API",
        "✅ Models with validation",
        "✅ 5 CRUD handlers with errors",
        "✅ Thread-safe repository",
        "✅ RESTful routes"
    ],
    "Act 1 tackles boilerplate. Traditional: 4 hours writing models, handlers, repository, routes. With Cloud Agent: 8 minutes from one prompt. 30x faster. Delivers complete enrollment API with models, 5 CRUD handlers, thread-safe repository, RESTful routes, comprehensive comments, and production-ready PR. Half a workday compressed into 8 minutes."
)

# SLIDE 7: Act 2 - Performance
add_stats_slide(
    prs,
    "⚡ ACT 2: PERFORMANCE AGENT",
    [
        "PROBLEM: 450ms response time (TOO SLOW!)",
        "Target: <100ms for production",
        "",
        "Traditional: 2.5 hours Redis integration",
        "AI Delegation: 6 minutes with Local Agent",
        "⚡ 25x FASTER",
        "",
        "RESULT: 12ms response time",
        "🚀 37x PERFORMANCE IMPROVEMENT",
        "✅ Cache hit/miss headers",
        "✅ Automatic invalidation"
    ],
    "Act 2: performance crisis. 450ms response time, need under 100ms. Traditional: 2.5 hours implementing Redis caching, connection pooling, invalidation logic. With Local Agent: 6 minutes from prompt. Result: 12ms response time, 37x improvement! We prove it live with before/after testing and Redis inspection."
)

# SLIDE 8: Act 3 - QA
add_stats_slide(
    prs,
    "🛡️ ACT 3: QA CONTRACT AGENT",
    [
        "CHALLENGE: Production reliability",
        "\"Fast code without quality is fast failure\"",
        "",
        "Traditional: 3 hours of QA work",
        "AI Delegation: 8 minutes with Local Agent",
        "⚡ 22x FASTER",
        "",
        "DELIVERED:",
        "✅ OpenAPI 3.0 specification",
        "✅ Contract validation script",
        "✅ Integration test suite",
        "✅ 94.2% test coverage"
    ],
    "Act 3: quality gates. Traditional: 3 hours writing OpenAPI specs, contract validation, integration tests. With Local Agent: 8 minutes. Delivers complete OpenAPI spec, contract validation preventing breaking changes, integration test suite with 94.2% coverage. We run all tests live to prove quality."
)

# SLIDE 9: Act 4 - Documentation
add_stats_slide(
    prs,
    "📚 ACT 4: DOCUMENTATION AGENT",
    [
        "REALITY: \"We'll document it later...\"",
        "(Later never comes)",
        "",
        "Traditional: 80 minutes of writing",
        "AI Delegation: 5 minutes Background Agent",
        "⚡ 16x FASTER",
        "",
        "DELIVERED:",
        "✅ Godoc comments (100% coverage)",
        "✅ Professional README",
        "✅ API usage examples",
        "✅ Troubleshooting guide",
        "",
        "💡 Works while we demo other acts!"
    ],
    "Act 4: documentation everyone hates. Traditional: 80 minutes writing godoc, README, examples. With Background Agent: 5 minutes running in parallel while we demo Acts 2 & 3. Delivers 100% API coverage, professional README, usage examples, troubleshooting guide. Documentation that people actually read and use."
)

# SLIDE 10: Transformation Summary
add_stats_slide(
    prs,
    "🏆 MISSION ACCOMPLISHED",
    [
        "FROM ZERO TO PRODUCTION IN 27 MINUTES",
        "",
        "✅ Complete CRUD API",
        "✅ Redis caching (37x faster)",
        "✅ OpenAPI spec + validation",
        "✅ 94.2% test coverage",
        "✅ 100% documentation",
        "",
        "Traditional: 10 hours • AI: 27 minutes",
        "⚡ 22x SPEEDUP",
        "",
        "Developer: $750 • Copilot: $0.05",
        "💰 ROI: 14,999%"
    ],
    "27 minutes of demo, 10 hours saved. Complete production-ready system. 22x speedup. Developer cost $750, Copilot cost 5 cents. That's 14,999% ROI. But more importantly: we maintained 94% coverage, contract validation, comprehensive docs. Speed WITHOUT sacrificing quality. Freed developers for creative work."
)

# SLIDE 11: Key Takeaways
add_content_slide(
    prs,
    "💡 KEY LESSONS",
    [
        "1️⃣  DELEGATE PATTERNS, NOT CREATIVITY",
        "  Boilerplate → AI delegation",
        "  Unique business logic → Human expertise",
        "",
        "2️⃣  CHOOSE THE RIGHT AGENT",
        "  Cloud → Complex features",
        "  Local → Interactive work",
        "  Background → Parallel tasks",
        "",
        "3️⃣  QUALITY & SPEED AREN'T OPPOSITES",
        "  AI generates tests faster than humans",
        "  Documentation happens automatically",
        "",
        "4️⃣  START SMALL, SCALE FAST",
        "  Pick ONE task tomorrow",
        "  Measure time saved",
        "  Success breeds adoption"
    ],
    "Four key lessons: 1) Delegate patterns to AI, keep creativity human. 2) Strategic agent choice matters more than speed. 3) AI helps you move fast AND maintain quality. 4) Start with one task, measure savings, let success spread naturally. Don't revolutionize tomorrow, evolve iteratively."
)

# SLIDE 12: Session 2 Preview
add_content_slide(
    prs,
    "🚀 COMING NEXT: SESSION 2",
    [
        "REFACTORING & QA EXPANSION",
        "Building on What AI Created",
        "",
        "🔧 REFACTORING WITH AI",
        "  • Transform concrete to abstractions",
        "  • Extract repository interfaces",
        "  • Apply design patterns automatically",
        "",
        "🧪 ADVANCED QA AUTOMATION",
        "  • Generate mock implementations",
        "  • Create test data factories",
        "  • Build comprehensive test suites",
        "",
        "💡 Evolution > Revolution",
        "Making AI-generated code enterprise-ready"
    ],
    "Session 2 preview: refactoring and advanced testing. Take AI-generated code and make it enterprise-ready. Transform concrete code to flexible abstractions. Generate mocks and test factories. Expand coverage from 94% to bulletproof. Evolution beats revolution - improve code iteratively with AI assistance."
)

# SLIDE 13: Q&A
add_content_slide(
    prs,
    "💬 QUESTIONS & NEXT STEPS",
    [
        "📋 YOUR CHALLENGE:",
        "",
        "Tomorrow morning:",
        "1️⃣  Pick ONE boring task",
        "2️⃣  Write ONE prompt to delegate it",
        "3️⃣  Measure time saved",
        "4️⃣  Calculate your ROI",
        "",
        "🔗 RESOURCES:",
        "  • Demo repository",
        "  • Session 1 slides",
        "  • Jira templates",
        "  • Agent prompt examples",
        "",
        "📅 UPCOMING SESSIONS:",
        "  Session 2: Refactoring & QA",
        "  Session 3: Security & DevOps",
        "  Session 4: Advanced Workflows"
    ],
    "Your homework: tomorrow, pick one boring task and delegate it. Measure time saved. Calculate ROI. Share results with your team. Resources available: demo repo, slides, Jira templates, prompt examples. Session 2 next: refactoring and advanced testing. Questions?"
)

# Save presentation
output_file = "SESSION1_GitHub_Copilot_Agent_Mode.pptx"
prs.save(output_file)
print(f"✅ PowerPoint created successfully: {output_file}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"🎯 Ready for your client presentation!")
